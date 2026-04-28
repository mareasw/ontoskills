"""GAIA benchmark wrapper.

Loads the GAIA dataset from HuggingFace, runs tasks through a BaseAgent
subclass, and scores results using exact-match comparison.

Dataset: ``gaia-benchmark/GAIA``
Levels:  ``2023_level1``, ``2023_level2``, ``2023_level3``
"""

from __future__ import annotations

import json
import logging
import os
import re
from pathlib import Path
from typing import Any

from datasets import load_dataset  # type: ignore[import-untyped]

from benchmark.agents.base import AgentResult, BaseAgent

logger = logging.getLogger(__name__)


def _normalize_extracted(val: str) -> str:
    """Strip units and trailing text from an extracted answer value."""
    val = val.strip()
    # If it starts with a number, strip trailing units/words
    m = re.match(r"^(-?[\d,]+(?:\.\d+)?)\s*(?:hours?|m²|m³|meters?|years?|days?|times?|dollars?|albums?|songs?|books?|people?|students?|items?|studio\s+\w+)?\.?$", val, re.I)
    if m:
        return m.group(1).replace(",", "")
    return val.rstrip(".")


def extract_answer(text: str) -> str:
    """Extract a short answer from a long-form model response.

    Tries several patterns in priority order:
    0. FINAL ANSWER: X (explicit format from system prompt)
    1. LaTeX \\boxed{...}
    2. Last bold-formatted value (``**...**``)
    3. "the answer is X" / "answer: X" patterns
    4. Last standalone number in the final sentence

    Returns the full text unchanged if no pattern matches.
    """
    if not text:
        return text

    # 0. FINAL ANSWER: X (highest priority)
    fa_matches = re.findall(r"FINAL\s+ANSWER:\s*(.+?)(?:\n|$)", text, re.IGNORECASE)
    if fa_matches:
        return _normalize_extracted(fa_matches[-1].strip())

    # 1. LaTeX \boxed{...}
    m = re.search(r"\\boxed\{([^}]+)\}", text)
    if m:
        return _normalize_extracted(m.group(1).strip())

    # 2. Last **bold** value that looks like a short answer
    bold_matches = re.findall(r"\*\*([^*]{1,80}?)\*\*", text)
    if bold_matches:
        _heading_words = ("answer", "solution", "result", "analysis", "note", "summary")
        for bm in reversed(bold_matches):
            bm_stripped = bm.strip()
            if "\n" in bm_stripped:
                continue
            if bm_stripped.lstrip("#").strip().lower() in _heading_words:
                continue
            if re.match(r"^[\d,.\-/+%°$€£]+$", bm_stripped):
                return _normalize_extracted(bm_stripped)
            if len(bm_stripped) <= 30 and not any(
                w in bm_stripped.lower() for w in ("the ", "is ", "are ", "this ", "we ", "look ", "use ")
            ):
                return _normalize_extracted(bm_stripped)

    # 3. "the answer is X" / "answer: X" patterns (not markdown headings)
    for line in text.split("\n"):
        line = line.strip()
        if line.startswith("#"):
            continue
        m = re.search(
            r"(?:the\s+(?:final\s+)?answer(?:\s+is)?|final\s+answer|answer:)\s*:?\s*(.+?)$",
            line,
            re.IGNORECASE,
        )
        if m:
            return _normalize_extracted(m.group(1).strip().rstrip("."))

    # 4. Last standalone number in the final line
    lines = [l.strip() for l in text.strip().split("\n") if l.strip()]
    if lines:
        last_line = lines[-1]
        nums = re.findall(r"(?<!\w)(-?[\d,]+(?:\.\d+)?)(?!\w)", last_line)
        if nums:
            return nums[-1].replace(",", "")

    return text


def _extract_final_answer(messages: list[dict]) -> str:
    """Extract the final text answer from conversation messages.

    Searches in reverse for the last assistant message containing text.
    """
    for block in reversed(messages):
        if isinstance(block, dict) and block.get("role") == "assistant":
            content = block.get("content", "")
            if isinstance(content, str):
                return content
            elif isinstance(content, list):
                texts = [
                    b["text"]
                    for b in content
                    if isinstance(b, dict) and b.get("type") == "text"
                ]
                answer = "\n".join(texts)
                if answer.strip():
                    return answer
    return ""


# Skill-relevant file types and question patterns for task filtering.
_SKILL_RELEVANT_EXTENSIONS = frozenset({
    ".py", ".xlsx", ".csv", ".docx", ".pptx", ".pdf", ".txt", ".json",
})
_SKILL_RELEVANT_PATTERNS = re.compile(
    r"\b(code|python|program|output|spreadsheet|excel|file|attached|attachment|"
    r"data|table|chart|formula|calculate|algorithm|function|class|method|"
    r"debug|error|exception|compile|execute|run|import|module|library|"
    r"api|sql|query|database|schema|json|csv|pandas|numpy|matplotlib|"
    r"machine learning|model|train|predict|neural|bert|transformer)\b",
    re.IGNORECASE,
)


def is_skill_relevant(task: dict) -> bool:
    """Check if a GAIA task is relevant to our skill set (file analysis, coding, data)."""
    file_path = task.get("file_path") or ""
    question = task.get("question", "")

    # Tasks with relevant file types are always included.
    if file_path:
        ext = Path(file_path).suffix.lower()
        if ext in _SKILL_RELEVANT_EXTENSIONS:
            return True
        # Skip audio/video files (we can't process them).
        if ext in (".mp3", ".wav", ".mp4", ".avi"):
            return False

    # Tasks with skill-relevant question content.
    if _SKILL_RELEVANT_PATTERNS.search(question):
        return True

    return False


# ---------------------------------------------------------------------------
# Tool schemas — given to both agents so they can read attachments & search
# ---------------------------------------------------------------------------

READ_FILE_TOOL: dict[str, Any] = {
    "name": "read_file",
    "description": "Read the contents of a file.",
    "input_schema": {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Absolute or relative path to the file to read.",
            },
        },
        "required": ["path"],
    },
}

WEB_SEARCH_TOOL: dict[str, Any] = {
    "name": "web_search",
    "description": "Search the web for information using Brave Search.",
    "input_schema": {
        "type": "object",
        "properties": {
            "query": {
                "type": "string",
                "description": "Search query string.",
            },
        },
        "required": ["query"],
    },
}

_BRAVE_API_KEY = os.environ.get("BRAVE_API_KEY", "")


def _call_brave_search(query: str) -> str:
    """Execute a Brave Search API call and return formatted results."""
    import requests

    try:
        resp = requests.get(
            "https://api.search.brave.com/res/v1/web/search",
            params={"q": query, "count": 5},
            headers={
                "Accept": "application/json",
                "Accept-Encoding": "gzip",
                "X-Subscription-Token": _BRAVE_API_KEY,
            },
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
    except Exception as exc:
        return f"Search error: {exc}"

    results = data.get("web", {}).get("results", [])
    if not results:
        return "No results found."

    parts: list[str] = []
    for i, r in enumerate(results[:5], 1):
        title = r.get("title", "")
        url = r.get("url", "")
        desc = r.get("description", "")
        parts.append(f"{i}. {title}\n   {url}\n   {desc}")
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Web page reader — fetches and extracts text from a URL
# ---------------------------------------------------------------------------

WEB_READ_TOOL: dict[str, Any] = {
    "name": "web_read",
    "description": "Read the text content of a web page by URL.",
    "input_schema": {
        "type": "object",
        "properties": {
            "url": {
                "type": "string",
                "description": "The URL of the web page to read.",
            },
        },
        "required": ["url"],
    },
}


def _call_web_read(url: str) -> str:
    """Fetch a web page and extract readable text."""
    import requests as _req
    try:
        resp = _req.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=15)
        resp.raise_for_status()
        html = resp.text
    except Exception as exc:
        return f"Error fetching {url}: {exc}"

    # Strip HTML tags — basic but sufficient for Wikipedia etc.
    import re as _re
    text = _re.sub(r"<script[^>]*>.*?</script>", "", html, flags=_re.DOTALL | _re.I)
    text = _re.sub(r"<style[^>]*>.*?</style>", "", text, flags=_re.DOTALL | _re.I)
    text = _re.sub(r"<[^>]+>", " ", text)
    text = _re.sub(r"\s+", " ", text).strip()
    # Truncate to ~8000 chars to avoid flooding context
    return text[:8000]


def _handle_read_file(file_path: str) -> tuple[str, bool]:
    """Read a file, handling PDF, DOCX, PPTX, XLSX, CSV and text. Returns (content, is_error)."""
    from pathlib import Path as _P

    p = _P(file_path)

    # YouTube URL detection
    if "youtube.com" in file_path or "youtu.be" in file_path:
        return _fetch_youtube_transcript(file_path), False

    if not p.exists():
        return f"Error: file not found: {file_path}", True

    suffix = p.suffix.lower()

    # PDF
    if suffix == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(str(p)) as pdf:
                pages = []
                for page in pdf.pages[:20]:  # max 20 pages
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                return "\n\n".join(pages)[:12000], False
        except Exception as exc:
            return f"Error reading PDF {file_path}: {exc}", True

    # DOCX
    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(p))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n".join(paragraphs)[:12000], False
        except Exception as exc:
            return f"Error reading DOCX {file_path}: {exc}", True

    # PPTX
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text.strip())
                if texts:
                    slides.append(f"Slide {i}: " + " | ".join(texts))
            return "\n".join(slides)[:12000], False
        except Exception as exc:
            return f"Error reading PPTX {file_path}: {exc}", True

    # XLSX / CSV
    if suffix in (".xlsx", ".csv"):
        try:
            import pandas as pd
            if suffix == ".csv":
                df = pd.read_csv(str(p))
            else:
                df = pd.read_excel(str(p))
            return df.to_string(max_rows=100, max_cols=20)[:12000], False
        except Exception as exc:
            return f"Error reading {suffix} {file_path}: {exc}", True

    # Text files
    try:
        return p.read_text(encoding="utf-8"), False
    except Exception as exc:
        return f"Error reading {file_path}: {exc}", True


def _fetch_youtube_transcript(url: str) -> str:
    """Extract transcript from a YouTube URL."""
    import re as _re
    from youtube_transcript_api import YouTubeTranscriptApi

    m = _re.search(
        r"(?:v=|youtu\.be/|embed/)([a-zA-Z0-9_-]{11})", url
    )
    if not m:
        return f"Could not extract YouTube video ID from: {url}"

    video_id = m.group(1)
    try:
        ytt = YouTubeTranscriptApi()
        transcript = ytt.fetch(video_id)
        lines = [entry.text if hasattr(entry, "text") else entry["text"] for entry in transcript]
        return f"YouTube transcript for {video_id}:\n" + " ".join(lines)[:8000]
    except Exception as exc:
        return f"Error fetching YouTube transcript for {video_id}: {exc}"


class GAIAWrapper:
    """GAIA benchmark wrapper.

    Parameters
    ----------
    data_dir:
        Directory to cache downloaded GAIA data and attachment files.
    """

    _VALID_LEVELS = ("2023_level1", "2023_level2", "2023_level3")
    _VALID_SPLITS = ("test", "validation")

    def __init__(self, data_dir: str = "benchmark/data/gaia") -> None:
        self.data_dir = Path(data_dir)
        self.data_dir.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------
    # Dataset loading
    # ------------------------------------------------------------------

    def load_dataset(
        self,
        level: str = "2023_level1",
        split: str = "test",
    ) -> list[dict]:
        """Load GAIA tasks from HuggingFace.

        Returns a list of dicts with keys:
        ``task_id``, ``question``, ``file_path``, ``gold_answer``.
        """
        if level not in self._VALID_LEVELS:
            raise ValueError(
                f"Invalid level {level!r}. Choose from {self._VALID_LEVELS}"
            )
        if split not in self._VALID_SPLITS:
            raise ValueError(
                f"Invalid split {split!r}. Choose from {self._VALID_SPLITS}"
            )

        # Use snapshot_download + local load to handle gated datasets.
        import os
        from huggingface_hub import snapshot_download

        token = os.environ.get("HF_TOKEN") or os.environ.get("HUGGING_FACE_HUB_TOKEN")
        if not token:
            raise ValueError(
                "HF_TOKEN or HUGGING_FACE_HUB_TOKEN environment variable required "
                "for GAIA dataset. Run: huggingface-cli login"
            )
        data_dir = snapshot_download(
            repo_id="gaia-benchmark/GAIA",
            repo_type="dataset",
            token=token,
        )
        ds = load_dataset(data_dir, level, split=split)

        tasks: list[dict] = []
        for row in ds:
            task_id = row["task_id"]
            # GAIA uses title-case column names
            question = row.get("Question", row.get("question", ""))

            # Handle file attachment.
            file_path: str | None = None
            raw_file = row.get("file_name") or row.get("file_path") or row.get("file")
            if raw_file:
                if isinstance(raw_file, dict) and "path" in raw_file:
                    attachment_dir = self.data_dir / "attachments" / level
                    attachment_dir.mkdir(parents=True, exist_ok=True)
                    dest = attachment_dir / Path(raw_file["path"]).name
                    if not dest.exists():
                        import shutil
                        shutil.copy2(raw_file["path"], dest)
                    file_path = str(dest)
                elif isinstance(raw_file, str) and raw_file.strip():
                    p = Path(raw_file)
                    if p.exists():
                        file_path = str(p)
                    else:
                        # Resolve relative to HF snapshot directory.
                        resolved = Path(data_dir) / raw_file
                        if resolved.exists():
                            file_path = str(resolved)
                        else:
                            file_path = raw_file

            # Gold answer (title-case "Final answer" in GAIA).
            gold_answer: str | None = row.get("Final answer") or row.get("final_answer", None)
            # Some gold answers are "?" (withheld).
            if gold_answer and gold_answer.strip() in ("?", ""):
                gold_answer = None

            tasks.append({
                "task_id": str(task_id),
                "question": question,
                "file_path": file_path,
                "gold_answer": gold_answer,
            })

        logger.info("Loaded %d GAIA tasks (level=%s, split=%s)", len(tasks), level, split)
        return tasks

    # ------------------------------------------------------------------
    # Single-task execution
    # ------------------------------------------------------------------

    def run_task(self, agent: BaseAgent, task: dict, *, mcp_client: Any = None) -> dict:
        """Run a single GAIA task through an agent.

        Both agents receive a ``read_file`` tool so they can inspect file
        attachments.  For TraditionalAgent (which normally has no tools) this
        is the only tool provided.

        If *mcp_client* is provided (already started), it is reused instead of
        starting a fresh MCP lifecycle per task.  This avoids ~8-10 s overhead
        on each task.

        Returns a dict with:
        ``task_id``, ``model_answer``, ``metrics`` (AgentResult).
        """
        # Build the task prompt, including file-attachment hint when present.
        prompt = task["question"]
        if task.get("file_path"):
            prompt += (
                f"\n\n[Attachment available at: {task['file_path']}. "
                "Use the read_file tool to inspect it if needed.]"
            )

        prompt += (
            "\n\nIMPORTANT: After your analysis, you MUST end your response "
            "with EXACTLY this format on its own line:\n"
            "FINAL ANSWER: <your answer>\n\n"
            "Your answer must be a single value — a number, word, or short "
            "phrase. Do NOT include units unless they are part of the answer. "
            "Do NOT include punctuation unless it is part of the answer. "
            "Do NOT add any text after the FINAL ANSWER line."
        )

        # Patch get_tools to include read_file.
        original_get_tools = agent.get_tools
        original_run_turn = agent.run_turn

        def _patched_get_tools() -> list[dict] | None:
            base_tools = original_get_tools()
            injected = [READ_FILE_TOOL, WEB_SEARCH_TOOL, WEB_READ_TOOL]
            if base_tools is None:
                return injected
            names = {t["name"] for t in base_tools}
            extra = [t for t in injected if t["name"] not in names]
            return [*base_tools, *extra] if extra else base_tools

        def _patched_run_turn(messages: list[dict]) -> tuple[dict, dict]:
            """Execute one turn with read_file handling.

            Routes read_file calls to local file I/O, delegates MCP tool
            names to the original agent for MCP routing, and handles all
            other tool calls by recording them.
            """
            import time as _time

            start = _time.perf_counter()
            response = agent._call_api(messages)
            latency_ms = (_time.perf_counter() - start) * 1000

            content_blocks: list[dict] = []
            for block in response.content:
                if block.type == "text":
                    content_blocks.append({
                        "type": "text",
                        "text": block.text,
                    })
                elif block.type == "tool_use":
                    content_blocks.append({
                        "type": "tool_use",
                        "id": block.id,
                        "name": block.name,
                        "input": block.input,
                    })

            assistant_msg: dict = {
                "role": "assistant",
                "content": content_blocks,
            }

            tool_calls = 0
            tool_result_blocks: list[dict] = []
            for block in content_blocks:
                if block.get("type") != "tool_use":
                    continue
                tool_calls += 1
                tool_name = block["name"]
                tool_input = block.get("input", {})

                if tool_name == "read_file":
                    file_path = tool_input.get("path", "")
                    result_text, is_error = _handle_read_file(file_path)
                elif tool_name == "web_search":
                    query = tool_input.get("query", "")
                    logger.info("Web search: %s", query[:80])
                    result_text = _call_brave_search(query)
                    is_error = False
                elif tool_name == "web_read":
                    url = tool_input.get("url", "")
                    logger.info("Web read: %s", url[:80])
                    result_text = _call_web_read(url)
                    is_error = False
                else:
                    # Delegate read_skill to the agent's skill resolver.
                    if tool_name == "read_skill" and hasattr(agent, "_resolve_skill"):
                        skill_name = tool_input.get("skill_name", "")
                        content = agent._resolve_skill(skill_name)
                        if content is None:
                            result_text = f"Skill not found: {skill_name}"
                            is_error = True
                        else:
                            result_text = content
                            is_error = False
                    else:
                        result_text = f"Error: unknown tool {tool_name}"
                        is_error = True

                tool_result_blocks.append({
                    "type": "tool_result",
                    "tool_use_id": block["id"],
                    "content": result_text,
                    "is_error": is_error,
                })

            if tool_result_blocks:
                messages.append(assistant_msg)
                messages.append({
                    "role": "user",
                    "content": tool_result_blocks,
                })

            metrics: dict = {
                "input_tokens": response.usage.input_tokens,
                "output_tokens": response.usage.output_tokens,
                "latency_ms": latency_ms,
                "tool_calls": tool_calls,
            }
            return assistant_msg, metrics

        agent.get_tools = _patched_get_tools  # type: ignore[assignment]
        agent.run_turn = _patched_run_turn  # type: ignore[assignment]

        try:
            # If no shared mcp_client was passed, start one locally.
            _mcp_started = False
            if mcp_client is None and hasattr(agent, "_mcp_client"):
                agent._mcp_client.__enter__()
                agent._mcp_client.initialize()
                _mcp_started = True

            # Pre-fetch skill knowledge if agent supports it and MCP is live.
            prefetched = ""
            if (
                mcp_client is not None
                and mcp_client._proc is not None
                and hasattr(agent, "prefetch_skills")
            ):
                try:
                    prefetched = agent.prefetch_skills(prompt)
                    if prefetched:
                        agent._prefetched_knowledge = prefetched
                        logger.info(
                            "Pre-fetched %d chars for task %s",
                            len(prefetched), task["task_id"],
                        )
                except Exception as exc:
                    logger.warning("Pre-fetch failed for %s: %s", task["task_id"], exc)

            # Custom run-loop (same pattern as SWE-bench to avoid
            # double-appending when run_turn also appends messages).
            messages: list[dict] = [{"role": "user", "content": prompt}]
            total_input = 0
            total_output = 0
            total_latency_ms = 0.0
            total_tool_calls = 0
            turns = 0

            for _ in range(25):
                assistant_msg, metrics = agent.run_turn(messages)
                turns += 1
                total_input += metrics["input_tokens"]
                total_output += metrics["output_tokens"]
                total_latency_ms += metrics["latency_ms"]
                total_tool_calls += metrics["tool_calls"]

                tool_use_blocks = [
                    b for b in (assistant_msg.get("content") or [])
                    if isinstance(b, dict) and b.get("type") == "tool_use"
                ]

                if tool_use_blocks:
                    # run_turn already appended assistant_msg + tool_results.
                    pass
                else:
                    messages.append(assistant_msg)
                    break

            # Extract final text answer.
            answer = _extract_final_answer(messages)

            # If empty (tool loop: all turns were tool calls with no text),
            # force one more turn asking for the answer.
            if not answer.strip():
                messages.append({
                    "role": "user",
                    "content": (
                        "Please provide your FINAL ANSWER now based on "
                        "all the information you have gathered. "
                        "Format: FINAL ANSWER: <your answer>"
                    ),
                })
                try:
                    forced_msg, forced_metrics = agent.run_turn(messages)
                    turns += 1
                    total_input += forced_metrics["input_tokens"]
                    total_output += forced_metrics["output_tokens"]
                    total_latency_ms += forced_metrics["latency_ms"]
                    total_tool_calls += forced_metrics["tool_calls"]
                except Exception:
                    pass
                answer = _extract_final_answer(messages)

            result = AgentResult(
                answer=answer,
                input_tokens=total_input,
                output_tokens=total_output,
                total_latency_ms=total_latency_ms,
                tool_calls=total_tool_calls,
                turns=turns,
            )
        except Exception as exc:
            logger.warning("Agent error on task %s: %s", task["task_id"], exc)
            result = AgentResult(
                answer=f"[Agent error: {exc}]",
                input_tokens=0,
                output_tokens=0,
                total_latency_ms=0.0,
                tool_calls=0,
                turns=0,
            )
        finally:
            agent.get_tools = original_get_tools  # type: ignore[assignment]
            agent.run_turn = original_run_turn  # type: ignore[assignment]
            if prefetched:
                agent._prefetched_knowledge = ""
            if _mcp_started:
                try:
                    agent._mcp_client.__exit__(None, None, None)
                except Exception:
                    pass

        return {
            "task_id": task["task_id"],
            "model_answer": result.answer,
            "metrics": result,
        }

    # ------------------------------------------------------------------
    # Full benchmark run
    # ------------------------------------------------------------------

    def run_benchmark(
        self,
        agent: BaseAgent,
        level: str = "2023_level1",
        split: str = "validation",
        max_tasks: int | None = None,
        shuffle: bool = True,
        seed: int = 42,
    ) -> list[dict]:
        """Run all (or *max_tasks*) GAIA tasks through *agent*.

        For OntoSkillsAgent, starts the MCP subprocess ONCE and reuses it
        across all tasks, avoiding ~8-10 s startup overhead per task.

        Returns a list of result dicts (one per task).
        """
        import random

        tasks = self.load_dataset(level=level, split=split)
        # Filter to only skill-relevant tasks.
        tasks = [t for t in tasks if is_skill_relevant(t)]
        logger.info("GAIA skill-relevant tasks: %d", len(tasks))
        if shuffle:
            random.Random(seed).shuffle(tasks)
        if max_tasks is not None:
            tasks = tasks[:max_tasks]

        # Start MCP once for all tasks (OntoSkillsAgent).
        mcp_client = None
        if hasattr(agent, "_mcp_client"):
            mcp_client = agent._mcp_client
            mcp_client.__enter__()
            mcp_client.initialize()

        results: list[dict] = []
        try:
            for i, task in enumerate(tasks, 1):
                logger.info("Task %d/%d: %s", i, len(tasks), task["task_id"])
                try:
                    result = self.run_task(agent, task, mcp_client=mcp_client)
                except Exception:
                    logger.exception("Task %s failed", task["task_id"])
                    result = {
                        "task_id": task["task_id"],
                        "model_answer": "",
                        "metrics": None,
                    }
                results.append(result)
        finally:
            if mcp_client is not None:
                try:
                    mcp_client.__exit__(None, None, None)
                except Exception:
                    pass

        return results

    # ------------------------------------------------------------------
    # Scoring
    # ------------------------------------------------------------------

    @staticmethod
    def score(results: list[dict], gold_answers: dict[str, str]) -> dict:
        """Score results against gold answers using case-insensitive exact match.

        Parameters
        ----------
        results:
            List of result dicts from ``run_benchmark`` / ``run_task``.
        gold_answers:
            Mapping of ``task_id`` -> gold answer string.

        Returns
        -------
        dict with ``accuracy`` (float 0-1) and ``per_task`` details.
        """
        per_task: list[dict] = []
        correct = 0
        total = 0

        for r in results:
            task_id = r["task_id"]
            model_answer = r.get("model_answer", "")
            gold = gold_answers.get(task_id)

            if gold is None:
                match = None  # no gold answer available
                extracted = None
            else:
                total += 1
                extracted = extract_answer(model_answer)
                match = extracted.strip().lower() == gold.strip().lower()
                if match:
                    correct += 1

            per_task.append({
                "task_id": task_id,
                "model_answer": model_answer,
                "extracted_answer": extracted,
                "gold_answer": gold,
                "correct": match,
            })

        accuracy = correct / total if total > 0 else 0.0
        return {
            "accuracy": accuracy,
            "correct": correct,
            "total": total,
            "per_task": per_task,
        }

    # ------------------------------------------------------------------
    # Submission output
    # ------------------------------------------------------------------

    def write_submission(self, results: list[dict], output_path: str) -> None:
        """Write results as JSONL (one JSON object per line).

        Each line: ``{"task_id": ..., "model_answer": ...}``
        """
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)

        with open(path, "w", encoding="utf-8") as f:
            for r in results:
                entry = {
                    "task_id": r["task_id"],
                    "model_answer": r.get("model_answer", ""),
                }
                f.write(json.dumps(entry, ensure_ascii=False) + "\n")

        logger.info("Wrote %d results to %s", len(results), path)
